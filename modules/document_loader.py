"""
Document Loader Module
Adapted from:
  - pdf-to-slides-ai-generator: PDFProcessor.extract_text_from_pdf()
  - presenton: DocumentsLoader (multi-format loading)

Supports PDF, TXT, DOCX, and PPTX document ingestion.
"""

import os
import tempfile
from typing import List


class DocumentLoader:
    """
    Loads and extracts text from uploaded documents.
    Supports PDF, TXT, DOCX, and PPTX file formats.
    """

    SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".docx", ".pptx"}

    def load_document(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()

        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {ext}. Supported: {self.SUPPORTED_EXTENSIONS}"
            )

        if ext == ".pdf":
            return self._extract_pdf(file_path)
        elif ext == ".txt":
            return self._extract_txt(file_path)
        elif ext == ".docx":
            return self._extract_docx(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)

    def load_from_bytes(self, content: bytes, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()

        if ext == ".pdf":
            return self._extract_pdf_from_bytes(content)
        elif ext == ".txt":
            return content.decode("utf-8", errors="ignore")
        elif ext == ".docx":
            return self._extract_docx_from_bytes(content)
        elif ext == ".pptx":
            return self._extract_pptx_from_bytes(content)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _extract_pdf(self, file_path: str) -> str:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()

    def _extract_pdf_from_bytes(self, content: bytes) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            return self._extract_pdf(tmp_path)
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    def _extract_txt(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from a DOCX file using python-docx."""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    # Preserve heading levels
                    if para.style and para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "")
                        try:
                            level_num = int(level)
                            paragraphs.append(f"{'#' * level_num} {para.text}")
                        except ValueError:
                            paragraphs.append(para.text)
                    else:
                        paragraphs.append(para.text)

            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)

            return "\n\n".join(paragraphs)
        except ImportError:
            raise ImportError("python-docx is required for DOCX support. Install with: pip install python-docx")

    def _extract_docx_from_bytes(self, content: bytes) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            return self._extract_docx(tmp_path)
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from an existing PPTX file."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_text.append(paragraph.text.strip())
                if slide_text:
                    text_parts.append(f"## Slide {slide_num}\n" + "\n".join(slide_text))
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Could not extract text from PPTX: {e}")

    def _extract_pptx_from_bytes(self, content: bytes) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            return self._extract_pptx(tmp_path)
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    def load_multiple(self, file_paths: List[str]) -> str:
        all_text = []
        for path in file_paths:
            text = self.load_document(path)
            if text:
                all_text.append(f"--- Document: {os.path.basename(path)} ---\n{text}")
        return "\n\n".join(all_text)
