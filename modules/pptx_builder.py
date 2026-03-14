"""
PPTX Builder Module
Adapted from:
  - pdf-to-slides-ai-generator: PPTGenerator (python-pptx slide creation by type)
  - presenton: PptxPresentationCreator (advanced slide styling, shapes, fonts)

Creates Stranger Things / Hawkins Lab themed PPTX presentations.
Uses full-slide rectangle shapes for dark backgrounds (more reliable
than slide.background.fill across all PowerPoint versions).
"""

import os
import uuid
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Theme Colors ──────────────────────────────────────────────────
HAWKINS_DARK = RGBColor(0x0B, 0x0C, 0x10)
HAWKINS_PANEL = RGBColor(0x14, 0x19, 0x22)
NEON_RED = RGBColor(0xFF, 0x2A, 0x2A)
CYAN_GLOW = RGBColor(0x00, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class PptxBuilder:
    """Generates themed PPTX presentations from slide outline JSON."""

    def __init__(self, output_dir: str = "slides"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def build(self, slide_data: Dict[str, Any]) -> str:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        slides = slide_data.get("slides", [])
        if not slides:
            slides = [{"type": "title", "title": slide_data.get("title", "Briefing"), "subtitle": slide_data.get("subtitle", "")}]

        for s in slides:
            t = s.get("type", "content")
            if t == "title":
                self._add_title_slide(prs, s)
            elif t == "bullet_points":
                self._add_bullet_slide(prs, s)
            elif t == "citation":
                self._add_citation_slide(prs, s)
            elif t == "section_divider":
                self._add_section_divider(prs, s)
            elif t == "two_column":
                self._add_two_column_slide(prs, s)
            elif t == "image_text":
                self._add_image_text_slide(prs, s)
            elif t == "chart_slide":
                self._add_chart_slide(prs, s)
            else:
                self._add_content_slide(prs, s)

        fid = str(uuid.uuid4())[:8]
        fp = os.path.join(self.output_dir, f"briefing_{fid}.pptx")
        prs.save(fp)
        return fp

    # ── helpers ───────────────────────────────────────────────────

    def _dark_bg(self, slide):
        """Add a full-slide dark rectangle as background (most reliable method)."""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = HAWKINS_DARK
        bg.line.fill.background()
        # Send to back
        sp = bg._element
        sp.getparent().insert(0, sp)

    def _accent_bar(self, slide):
        """Cyan left accent bar."""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.3), Inches(0.08), Inches(6.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = CYAN_GLOW
        bar.line.fill.background()

    def _title_bar(self, slide):
        """Dark panel behind the title."""
        tb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
        tb.fill.solid()
        tb.fill.fore_color.rgb = HAWKINS_PANEL
        tb.line.fill.background()

    def _add_title_text(self, slide, text, top=Inches(0.2)):
        box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.9))
        p = box.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = NEON_RED
        p.font.name = "Arial Black"

    def _notes(self, slide, data):
        n = data.get("speaker_notes", "")
        if n:
            slide.notes_slide.notes_text_frame.text = str(n)

    # ── Title Slide ───────────────────────────────────────────────

    def _add_title_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._dark_bg(slide)

        # Red accent line top
        lt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.3), Pt(3))
        lt.fill.solid(); lt.fill.fore_color.rgb = NEON_RED; lt.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(data.get("title", "CLASSIFIED BRIEFING"))
        p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = NEON_RED
        p.font.name = "Arial Black"; p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sb = slide.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(10.3), Inches(1))
        sf = sb.text_frame; sf.word_wrap = True
        sp = sf.paragraphs[0]
        sp.text = str(data.get("subtitle", data.get("content", "")))
        sp.font.size = Pt(22); sp.font.color.rgb = CYAN_GLOW
        sp.font.name = "Arial"; sp.alignment = PP_ALIGN.CENTER

        # Red accent line bottom
        lb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.8), Inches(10.3), Pt(3))
        lb.fill.solid(); lb.fill.fore_color.rgb = NEON_RED; lb.line.fill.background()

        # Classification badge
        bb = slide.shapes.add_textbox(Inches(3.5), Inches(6.4), Inches(6.3), Inches(0.5))
        bp = bb.text_frame.paragraphs[0]
        bp.text = "HAWKINS NATIONAL LABORATORY — CLASSIFIED"
        bp.font.size = Pt(11); bp.font.color.rgb = NEON_RED
        bp.font.name = "Courier New"; bp.alignment = PP_ALIGN.CENTER

        self._notes(slide, data)

    # ── Bullet Slide ──────────────────────────────────────────────

    def _add_bullet_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._dark_bg(slide)
        self._title_bar(slide)
        self._accent_bar(slide)
        self._add_title_text(slide, data.get("title", "Intel Report"))

        bullets = data.get("bullet_points", [])
        if bullets:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11), Inches(5))
            tf = box.text_frame; tf.word_wrap = True
            for i, bt in enumerate(bullets[:6]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = "▸  " + str(bt)
                p.font.size = Pt(20); p.font.color.rgb = LIGHT_GRAY
                p.font.name = "Consolas"; p.space_after = Pt(18)

        self._notes(slide, data)

    # ── Content Slide ─────────────────────────────────────────────

    def _add_content_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._dark_bg(slide)
        self._title_bar(slide)
        self._accent_bar(slide)
        self._add_title_text(slide, data.get("title", "Analysis"))

        content = data.get("content", "")
        if not content and data.get("bullet_points"):
            content = "\n\n".join(str(b) for b in data["bullet_points"])
        if content:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11), Inches(5))
            tf = box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(content)
            p.font.size = Pt(20); p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Consolas"; p.line_spacing = Pt(30)

        self._notes(slide, data)

    # ── Section Divider Slide ────────────────────────────────────

    def _add_section_divider(self, prs, data):
        """Full-screen centered heading to separate major sections."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._dark_bg(slide)

        # Top accent line
        lt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2.5), Inches(7.3), Pt(3))
        lt.fill.solid(); lt.fill.fore_color.rgb = NEON_RED; lt.line.fill.background()

        # Section title — large centered
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1.5))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(data.get("title", "NEXT SECTION"))
        p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = NEON_RED
        p.font.name = "Arial Black"; p.alignment = PP_ALIGN.CENTER

        # Subtitle — cyan below
        sub = data.get("subtitle", "")
        if sub:
            sb = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.3), Inches(0.8))
            sf = sb.text_frame; sf.word_wrap = True
            sp = sf.paragraphs[0]
            sp.text = str(sub)
            sp.font.size = Pt(20); sp.font.color.rgb = CYAN_GLOW
            sp.font.name = "Arial"; sp.alignment = PP_ALIGN.CENTER

        # Bottom accent line
        lb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(5.0), Inches(7.3), Pt(3))
        lb.fill.solid(); lb.fill.fore_color.rgb = NEON_RED; lb.line.fill.background()

        self._notes(slide, data)

    # ── Two-Column Slide ─────────────────────────────────────────

    def _add_two_column_slide(self, prs, data):
        """Side-by-side bullet lists for comparisons."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._dark_bg(slide)
        self._title_bar(slide)
        self._accent_bar(slide)
        self._add_title_text(slide, data.get("title", "Comparison"))

        left_col = data.get("left_column", [])
        right_col = data.get("right_column", [])

        # Vertical divider line
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.65), Inches(1.8), Pt(2), Inches(5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x44)
        divider.line.fill.background()

        # Left column
        if left_col:
            lbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
            ltf = lbox.text_frame; ltf.word_wrap = True
            for i, bt in enumerate(left_col[:6]):
                p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
                p.text = "▸  " + str(bt)
                p.font.size = Pt(18); p.font.color.rgb = LIGHT_GRAY
                p.font.name = "Consolas"; p.space_after = Pt(14)

        # Right column
        if right_col:
            rbox = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(5))
            rtf = rbox.text_frame; rtf.word_wrap = True
            for i, bt in enumerate(right_col[:6]):
                p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
                p.text = "▸  " + str(bt)
                p.font.size = Pt(18); p.font.color.rgb = LIGHT_GRAY
                p.font.name = "Consolas"; p.space_after = Pt(14)

        self._notes(slide, data)

    # ── Image + Text Slide ───────────────────────────────────────

    def _add_image_text_slide(self, prs, data):
        """Title + image (real or placeholder) + text content."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._dark_bg(slide)
        self._title_bar(slide)
        self._accent_bar(slide)
        self._add_title_text(slide, data.get("title", "Visual Analysis"))

        image_path = data.get("image_path", "")
        has_image = image_path and os.path.exists(image_path)

        if has_image:
            # Insert real image with aspect-ratio preservation
            try:
                from PIL import Image as PILImage
                with PILImage.open(image_path) as img:
                    img_w, img_h = img.size
            except Exception:
                img_w, img_h = 800, 600

            max_w = Inches(5.5)
            max_h = Inches(4.5)

            aspect = img_w / img_h if img_h else 1
            if aspect > (5.5 / 4.5):
                w = max_w
                h = Emu(int(max_w / aspect))
            else:
                h = max_h
                w = Emu(int(max_h * aspect))

            left = Inches(0.8) + (max_w - w) // 2
            top = Inches(1.8) + (max_h - h) // 2

            try:
                slide.shapes.add_picture(image_path, left, top, w, h)
            except Exception as e:
                print(f"[PptxBuilder] Failed to insert image: {e}")
                has_image = False

        if not has_image:
            # Fallback: dark placeholder panel
            img_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5)
            )
            img_box.fill.solid()
            img_box.fill.fore_color.rgb = HAWKINS_PANEL
            img_box.line.color.rgb = RGBColor(0x33, 0x44, 0x55)
            img_box.line.width = Pt(1)

            label = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(4), Inches(0.6))
            lp = label.text_frame.paragraphs[0]
            desc = data.get("image_description", "[Visual content]")
            lp.text = str(desc) if desc else "[Visual content]"
            lp.font.size = Pt(14); lp.font.color.rgb = CYAN_GLOW
            lp.font.name = "Courier New"; lp.alignment = PP_ALIGN.CENTER

        # Right-side text content
        content = data.get("content", "")
        if content:
            tbox = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(5))
            tf = tbox.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(content)
            p.font.size = Pt(18); p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Consolas"; p.line_spacing = Pt(28)

        self._notes(slide, data)

    # ── Citation Slide ────────────────────────────────────────────

    def _add_citation_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._dark_bg(slide)
        self._title_bar(slide)
        self._accent_bar(slide)
        self._add_title_text(slide, data.get("title", "References"))

        refs = data.get("content", data.get("references", ""))
        if isinstance(refs, list):
            refs = "\n".join(str(r) for r in refs)
        if refs:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11), Inches(5))
            tf = box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(refs)
            p.font.size = Pt(16); p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Courier New"

        self._notes(slide, data)

    # ── Chart Slide ─────────────────────────────────────────────

    def _add_chart_slide(self, prs, data):
        """Title + large centered chart image."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._dark_bg(slide)
        self._title_bar(slide)
        self._accent_bar(slide)
        self._add_title_text(slide, data.get("title", "Data Visualization"))

        chart_path = data.get("chart_path", "")
        if chart_path and os.path.exists(chart_path):
            # Center large chart in the slide
            max_w = Inches(11.0)
            max_h = Inches(5.5)
            
            # Slide is 13.33 x 7.5. Center horizontally.
            left = (SLIDE_W - max_w) // 2
            top = Inches(1.8)
            
            try:
                slide.shapes.add_picture(chart_path, left, top, max_w, max_h)
            except Exception as e:
                print(f"[PptxBuilder] Failed to insert chart: {e}")
        
        self._notes(slide, data)
